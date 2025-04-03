#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import io
from pptx import Presentation
import json

# Set stdout to use UTF-8 encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    
    for slide_number, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Ensure text is properly encoded
                    text = text.encode('utf-8', errors='ignore').decode('utf-8')
                    slide_text.append(text)
        
        # Extract text from tables
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            # Ensure text is properly encoded
                            cell_text = cell_text.encode('utf-8', errors='ignore').decode('utf-8')
                            row_text.append(cell_text)
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        
        if slide_text:
            text_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))
    
    return "\n\n".join(text_content)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python parse_pptx.py <file_path>", file=sys.stderr)
        sys.exit(1)
    
    file_path = sys.argv[1]
    try:
        content = extract_text_from_pptx(file_path)
        print(content)
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1) 